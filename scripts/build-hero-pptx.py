#!/usr/bin/env python3
"""Assemble the hero layers into a 1920x1080 PPTX (one slide, every element a
separate draggable image) for import into Canva as an editable design."""
from pptx import Presentation
from pptx.util import Emu
from pptx.dml.color import RGBColor

PX = 9525  # EMU per pixel at 96dpi
L = "assets/hero-layers"

prs = Presentation()
prs.slide_width = Emu(1920 * PX)
prs.slide_height = Emu(1080 * PX)
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

bg = slide.background.fill
bg.solid()
bg.fore_color.rgb = RGBColor(0xF9, 0xE7, 0xCB)

# (file, left, top, width, height, rotation)
items = [
    (f"{L}/layer-word-just.png",   84, 120, 720, 300, 0),
    (f"{L}/layer-word-like.png",  777, 400, 600, 300, 0),
    (f"{L}/layer-word-that.png",   84, 756, 770, 300, 0),
    (f"{L}/layer-swoosh-arrow.png", 980, 330, 480, 260, 0),
    (f"{L}/layer-snap-ticks.png", 790, 650, 160, 150, 0),
    (f"{L}/layer-milo-engraved.png", 1254, 505, 630, 625, 0),
    (f"{L}/layer-sparkle.png",    950, 170,  56,  56, 0),
    (f"{L}/layer-sparkle.png",    726, 480,  40,  40, 0),
    (f"{L}/layer-sparkle.png",    586,  54,  36,  36, 0),
    (f"{L}/layer-photo-chip.png", 800,  50, 330, 350, 6),
    (f"{L}/layer-plaque-engraved.png", 860, 735, 499, 329, -5),
    (f"{L}/layer-eyebrow.png",     84,  56, 480,  40, 0),
    ("brand/logo/logo-full-250.png", 1758, 48, 78, 40, 0),
    (f"{L}/layer-tagline.png",   1430, 142, 440,  50, 0),
    (f"{L}/layer-cta-pill.png",  1554, 190, 330, 100, 0),
]
for path, left, top, w, h, rot in items:
    pic = slide.shapes.add_picture(path, Emu(left*PX), Emu(top*PX), Emu(w*PX), Emu(h*PX))
    if rot:
        pic.rotation = rot

prs.save(f"{L}/hero-layers.pptx")
print("saved", f"{L}/hero-layers.pptx")
