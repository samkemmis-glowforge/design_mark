#!/usr/bin/env python3
"""Assemble the triptych layers into a 1920x1080 PPTX (every element its own
draggable object) for Canva import. Reads geometry from
assets/hero-layers/triptych-layers.json (written by export-triptych-layers.ts)."""
import json

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Emu

PX = 9525
A, L = "assets/magic-engraver", "assets/hero-layers"
g = json.load(open(f"{L}/triptych-layers.json"))
W, H = g["canvas"]
P = g["panel"]

def hexc(s): return RGBColor(*(int(s[i:i+2], 16) for i in (1, 3, 5)))

# helper PNGs: cream word recolors + P3 scene crop
for wspec in g["words"]:
    if wspec["color"] != "cream":
        continue
    im = Image.open(f"{L}/{wspec['src']}").convert("RGBA")
    px = im.load()
    cream = (0xF9, 0xE7, 0xCB)
    for yy in range(im.height):
        for xx in range(im.width):
            r, gg, b, al = px[xx, yy]
            if al:
                px[xx, yy] = (*cream, al)
    out = wspec["src"].replace(".png", "-cream.png")
    im.save(f"{L}/{out}")
    wspec["src"] = out

scene = Image.open(f"{A}/coaster-scene.png").convert("RGB")
sw = round(g["scene"]["sceneW"])
scene = scene.resize((sw, sw))
ox = g["scene"]["panelX"] - g["scene"]["sceneLeft"]  # crop offset into scene
scene.crop((round(ox), 0, round(ox) + P, H)).save(f"{L}/triptych-p3-scene.png")

prs = Presentation()
prs.slide_width, prs.slide_height = Emu(W * PX), Emu(H * PX)
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide.background.fill.solid()
slide.background.fill.fore_color.rgb = hexc(g["bg"])

def rect(x, y, w, h, color):
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Emu(x*PX), Emu(y*PX), Emu(w*PX), Emu(h*PX))
    sh.fill.solid(); sh.fill.fore_color.rgb = hexc(color); sh.line.fill.background()

def pic(path, p):
    slide.shapes.add_picture(path, Emu(round(p["x"])*PX), Emu(round(p["y"])*PX),
                             Emu(round(p["w"])*PX), Emu(round(p["h"])*PX))

rect(P, 0, P, H, g["p2"])                                   # teal panel
pic(f"{L}/triptych-p3-scene.png", {"x": 2*P, "y": 0, "w": P, "h": H})
rect(P - 2, 0, 3, H, g["ink"]); rect(2*P - 2, 0, 3, H, g["ink"])  # dividers
for wspec in g["words"]:
    pic(f"{L}/{wspec['src']}", wspec)
pic(f"{L}/triptych-photo-ring.png", g["photoRing"])
pic(f"{A}/milo2-cutout.png", g["photoPop"])
pic(f"{A}/m2-cut-teal.png", g["dogTeal"])
pic(f"{A}/engrave2-burnt.png", g["engraving"])
pic(f"{A}/name-milo.png", g["name"])

prs.save(f"{L}/triptych-layers.pptx")
print("saved", f"{L}/triptych-layers.pptx")
